"""JSON 规格 → pptx。用法：python make_pptx.py <spec.json> <out.pptx>"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

PRIMARY = RGBColor(0x1F, 0x4E, 0x79)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_section_slide(prs, sec):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = sec.get("heading", "")
    top = Inches(1.5)
    bullets = sec.get("bullets") or []
    if bullets:
        box = slide.shapes.add_textbox(Inches(0.6), top, Inches(9), Inches(0.4 * len(bullets) + 0.3))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(18)
        top = top + Inches(0.4 * len(bullets) + 0.5)
    table_spec = sec.get("table")
    if table_spec:
        headers = table_spec.get("headers") or []
        rows = table_spec.get("rows") or []
        if headers and rows is not None:
            shape = slide.shapes.add_table(
                len(rows) + 1, len(headers), Inches(0.6), top, Inches(9),
                Inches(0.35 * (len(rows) + 1)))
            table = shape.table
            for c, h in enumerate(headers):
                cell = table.cell(0, c)
                cell.text = str(h)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(14)
                    p.font.bold = True
            for r, row in enumerate(rows, start=1):
                for c, val in enumerate(row[: len(headers)]):
                    cell = table.cell(r, c)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(13)


def main():
    spec_path, out_path = sys.argv[1], sys.argv[2]
    spec = json.loads(Path(spec_path).read_text())
    prs = Presentation()
    add_title_slide(prs, spec.get("title", "报告"), spec.get("subtitle", ""))
    for sec in spec.get("sections", []):
        add_section_slide(prs, sec)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"OK {out_path} {len(prs.slides._sldIdLst)}页")


if __name__ == "__main__":
    main()
